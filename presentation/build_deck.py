#!/usr/bin/env python3
"""Build the STP final-presentation deck: presentation/STP-final-presentation.pptx.

Nomura Tech Graduate Program 2026 — final team presentation (18 slides + 3 appendix).
16:9, native editable PowerPoint shapes only (no images, no baked-in text).

Content source of truth: presentation/script.md (flow, quotes, metrics, honesty map).
Nothing in this deck may assert anything the script/repo does not back up.

Dependency: python-pptx — presentation tooling only, deliberately NOT pinned in
backend/requirements.txt (the app never imports it). Install into the project venv:
    backend/.venv/Scripts/pip install python-pptx

Regenerate (from the repo root):
    backend/.venv/Scripts/python presentation/build_deck.py

The script re-opens the generated file and verifies: 21 slides, non-empty speaker
notes on every slide, no empty text frames on content-sized shapes, and no
off-canvas shapes.

Edit guide:
- Wording of a slide      -> the slide_XX functions below (one per slide).
- Diagram layout          -> same functions; everything is positioned shapes.
- Speaker notes           -> the notes(...) call at the end of each slide function.
- Palette / fonts / chrome -> the constants and helpers in this header.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
OUT = HERE / "STP-final-presentation.pptx"

# ---------------------------------------------------------------- palette ----
# Conservative financial-industry palette: white bg, navy titles, gray-blue body,
# ONE accent (the product's #2962FF), muted gray for secondary text.
NAVY = RGBColor(0x1F, 0x38, 0x64)
BODY = RGBColor(0x33, 0x3F, 0x50)
ACCENT = RGBColor(0x29, 0x62, 0xFF)
MUTED = RGBColor(0x8A, 0x94, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF4, 0xF6, 0xFA)        # very light gray card fill
TINT = RGBColor(0xE9, 0xEF, 0xFC)         # very light accent tint (callouts)
HAIR = RGBColor(0xD9, 0xDE, 0xE7)         # hairline rules / card borders

FONT = "Segoe UI"

# 16:9 at the standard EMU size (13.333 x 7.5 in).
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN          # usable width between side margins


# ---------------------------------------------------------------- helpers ----
def R(t: str, s: float = 13, c: RGBColor = BODY, b: bool = False, i: bool = False):
    """Run spec: text, size(pt), color, bold, italic."""
    return {"t": t, "s": s, "c": c, "b": b, "i": i}


def _fill_runs(p, runs, space_after, line_spacing):
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    for r in runs:
        run = p.add_run()
        run.text = r["t"]
        f = run.font
        f.name = FONT
        f.size = Pt(r["s"])
        f.color.rgb = r["c"]
        f.bold = r["b"]
        f.italic = r["i"]


def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=4, line_spacing=1.0):
    """Plain text box. `lines` = list of paragraphs; each paragraph = list of R()."""
    tb = slide.shapes.add_textbox(_E(x), _E(y), _E(w), _E(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, runs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _fill_runs(p, runs, space_after, line_spacing)
    return tb


def card(slide, x, y, w, h, fill=PANEL, line=HAIR, radius=0.10,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Rounded-rectangle card (native shape, editable)."""
    sh = slide.shapes.add_shape(shape, _E(x), _E(y), _E(w), _E(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    _no_shadow(sh)
    return sh


def _E(v):
    """Coerce a Length/float to integer EMU — divisions like W/2 yield floats,
    which python-pptx writes as 'x.0' and then cannot read back."""
    return Emu(int(round(v)))


def shape_text(sh, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
               space_after=2, line_spacing=1.0, m=0.12):
    """Put formatted paragraphs into an autoshape's text frame."""
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(m)
    tf.margin_top = tf.margin_bottom = Inches(max(m - 0.04, 0.02))
    for i, runs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _fill_runs(p, runs, space_after, line_spacing)
    return sh


def add_bullets(slide, x, y, w, h, items, size=12.5, space_after=6, line_spacing=1.04):
    """Bulleted list. Item = (lead, rest, level):
    lead  -> bold navy lead-in run (or the whole bullet if rest is empty)
    rest  -> normal body run appended after the lead
    level -> 0 = '•' accent bullet, 1 = indented '–' muted sub-bullet
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (lead, rest, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        p.level = level
        glyph, gc = ("•  ", ACCENT) if level == 0 else ("–  ", MUTED)
        r = p.add_run()
        r.text = glyph
        r.font.name, r.font.size, r.font.color.rgb, r.font.bold = FONT, Pt(size), gc, True
        for text, color, bold in ((lead, NAVY if rest else BODY, bool(rest)),
                                  (rest, BODY, False)):
            if not text:
                continue
            r = p.add_run()
            r.text = text
            r.font.name, r.font.size = FONT, Pt(size)
            r.font.color.rgb, r.font.bold = color, bold
    return tb


def connect(slide, x1, y1, x2, y2, color=MUTED, w=1.5, arrow=True, dash=False):
    """Straight connector with an optional arrowhead (native, editable)."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _E(x1), _E(y1),
                                      _E(x2), _E(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    _no_shadow(conn)
    ln = conn.line._get_or_add_ln()
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    if arrow:
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    return conn


def base_slide(prs, idx, title, kicker=None):
    """Blank slide with the standard chrome: top accent bar, kicker, title,
    hairline rule under the title, slide number bottom-right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.055))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    _no_shadow(bar)
    if kicker:
        add_text(slide, MARGIN, Inches(0.30), CONTENT_W, Inches(0.24),
                 [[R(kicker.upper(), 10.5, MUTED, True)]])
    add_text(slide, MARGIN, Inches(0.52), CONTENT_W, Inches(0.62),
             [[R(title, 28, NAVY, True)]])
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.20),
                                  CONTENT_W, Pt(1.0))
    rule.fill.solid()
    rule.fill.fore_color.rgb = HAIR
    rule.line.fill.background()
    _no_shadow(rule)
    add_text(slide, SLIDE_W - MARGIN - Inches(0.6), SLIDE_H - Inches(0.40),
             Inches(0.6), Inches(0.24), [[R(str(idx), 10, MUTED)]],
             align=PP_ALIGN.RIGHT)
    return slide


def strip(slide, y, h, lines, fill=TINT, align=PP_ALIGN.CENTER, m=0.16):
    """Full-content-width callout strip near the bottom of a slide."""
    sh = card(slide, MARGIN, y, CONTENT_W, h, fill=fill, line=None)
    shape_text(sh, lines, align=align, m=m)
    return sh


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ============================================================== slides 1–7 ====
def slide_01(prs):
    """Title slide (custom chrome, no standard header)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.09))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    _no_shadow(bar)
    add_text(s, MARGIN, Inches(1.55), CONTENT_W, Inches(0.3),
             [[R("NOMURA TECH GRADUATE PROGRAM 2026  ·  FINAL TEAM PRESENTATION",
                 12, MUTED, True)]], align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(2.15), CONTENT_W, Inches(1.6),
             [[R("Next-Generation Trading Platform", 34, NAVY, True)],
              [R("with Straight-Through Processing", 34, NAVY, True)]],
             align=PP_ALIGN.CENTER, space_after=2)
    dash = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              _E(SLIDE_W / 2 - Inches(0.6)), Inches(3.72),
                              Inches(1.2), Inches(0.045))
    dash.fill.solid()
    dash.fill.fore_color.rgb = ACCENT
    dash.line.fill.background()
    _no_shadow(dash)
    add_text(s, MARGIN, Inches(4.0), CONTENT_W, Inches(0.4),
             [[R("An order goes from ticket to settlement with zero manual steps.",
                 15, BODY, False, True)]], align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(5.15), CONTENT_W, Inches(0.35),
             [[R("[Team name]  —  [Member 1] · [Member 2] · [Member 3] · "
                 "[Member 4]", 12.5, NAVY, True)]],
             align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(5.6), CONTENT_W, Inches(0.3),
             [[R("30 July 2026", 12, MUTED)]], align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(6.7), CONTENT_W, Inches(0.3),
             [[R("Training environment — market data is a replayed simulation dataset",
                 10, MUTED, False, True)]], align=PP_ALIGN.CENTER)
    notes(s, """
Good morning. We're [team name], and in three weeks we built a working trading
platform with straight-through processing — an order goes from ticket to
settlement with zero manual steps. But we didn't start with technology. We
started with four people.

(Replace [Team name] / [Member …] placeholders; presenters P1–P4 per script §1:
P1 corporate lead, P2/P4 technology, P3 corporate.)
""")


def slide_02(prs):
    s = base_slide(prs, 2, "Agenda", kicker="10 minutes + 5 minutes Q&A")
    items = [
        ("1", "The ask — four stakeholder voices", "P1 · 0:00–0:45"),
        ("2", "What we built + live demo", "P2 · 0:45–4:15"),
        ("3", "Operational processes", "P3 · 4:15–5:45"),
        ("4", "How we worked", "P4 · 5:45–7:45"),
        ("5", "What worked / what didn't", "P1 · 7:45–9:00"),
        ("6", "Learnings, roadmap & close", "P3 · 9:00–10:00"),
        ("7", "Q&A", "all · 10:00–15:00"),
    ]
    y = Inches(1.55)
    for num, title, who in items:
        chip = card(s, MARGIN, y, Inches(0.42), Inches(0.42), fill=NAVY,
                    line=None, radius=0.25)
        shape_text(chip, [[R(num, 14, WHITE, True)]], align=PP_ALIGN.CENTER, m=0.02)
        add_text(s, Inches(1.3), y + Inches(0.03), Inches(8.2), Inches(0.38),
                 [[R(title, 14.5, NAVY, True)]])
        add_text(s, Inches(9.6), y + Inches(0.06), Inches(3.1), Inches(0.34),
                 [[R(who, 11, MUTED)]], align=PP_ALIGN.RIGHT)
        y += Inches(0.62)
    notes(s, """
Quick roadmap of the next ten minutes: what we built — including a live
90-second demo — how it supports real operational workflows, how we worked as
a cross-functional team, what held up and what didn't, and where this goes
next.

Timing per script §1: the demo section gets the biggest block (35%). Appendix
slides A1–A3 stay on standby for Q&A.
""")


def slide_03(prs):
    s = base_slide(prs, 3, "The ask — four stakeholder voices", kicker="Section 1 · Hook")
    quotes = [
        ("Rohan", "Head of Product Development",
         "\u201cI want single-click trading.\u201d"),
        ("Tom & Patricia", "our clients",
         "\u201cThe current way is fine — don't make me learn another system. "
         "And stop sending me Excel.\u201d"),
        ("Nora", "the developer who maintains today's batch system",
         "\u201cWhy rebuild what we already have? Real-time is doable, "
         "but challenging.\u201d"),
        ("Roy", "our CTO",
         "\u201cSecurity, operations, SRE — DevSecOps from day one, "
         "or don't build it at all.\u201d"),
    ]
    W, H = Inches(5.95), Inches(2.32)
    for i, (name, role, quote) in enumerate(quotes):
        x = MARGIN + (i % 2) * (W + Inches(0.23))
        y = Inches(1.5) + (i // 2) * (H + Inches(0.2))
        c = card(s, x, y, W, H, fill=PANEL, line=HAIR)
        shape_text(c, [
            [R(name + "  ", 14, NAVY, True), R("— " + role, 10.5, MUTED)],
            [R(quote, 13.5, BODY, False, True)],
        ], anchor=MSO_ANCHOR.MIDDLE, space_after=8, m=0.22)
    strip(s, Inches(6.5), Inches(0.55), [
        [R("Everything you're about to see is our answer to those four sentences "
           "— hold us to it at the end.", 12.5, NAVY, True, True)]])
    notes(s, """
Rohan (Head of Product Development): "I want single-click trading."
Tom & Patricia (clients): "The current way is fine — don't make me learn
another system. And stop sending me Excel."
Nora (maintains today's batch system): "Why rebuild what we already have?
Real-time is doable, but challenging."
Roy (CTO): "Security, operations, SRE — DevSecOps from day one, or don't build
it at all."

Everything you're about to see is our answer to those four sentences. Hold us
to it at the end.

Handover: "To show you what that answer looks like — [P2]."
""")


def slide_04(prs):
    s = base_slide(prs, 4, "Goals & success criteria", kicker="Section 2 · What we built")
    add_text(s, MARGIN, Inches(1.42), Inches(5.9), Inches(0.3),
             [[R("THE BRIEF — FIVE CORE MODULES", 10.5, MUTED, True)]])
    modules = ["Order Execution", "Portfolio Management", "Reporting & Charting",
               "Technical Analytics", "Paper Trading"]
    y = Inches(1.82)
    for mname in modules:
        chip = card(s, MARGIN, y, Inches(5.9), Inches(0.72), fill=PANEL, line=HAIR)
        shape_text(chip, [[R(mname, 13.5, NAVY, True)]], m=0.2)
        y += Inches(0.86)
    x2 = Inches(6.85)
    w2 = Inches(5.88)
    c1 = card(s, x2, Inches(1.82), w2, Inches(1.5), fill=PANEL, line=HAIR)
    shape_text(c1, [
        [R("On a DevSecOps foundation", 13.5, NAVY, True)],
        [R("SSO-style auth, deny-by-default RBAC, tamper-evident audit, "
           "secrets behind a provider abstraction.", 11.5, BODY)],
    ], space_after=5, m=0.2)
    c2 = card(s, x2, Inches(3.52), w2, Inches(1.5), fill=PANEL, line=HAIR)
    shape_text(c2, [
        [R("Using the program's simulation dataset", 13.5, NAVY, True)],
        [R("Replayed market data — nothing in the system knows "
           "the data is simulated.", 11.5, BODY)],
    ], space_after=5, m=0.2)
    c3 = card(s, x2, Inches(5.22), w2, Inches(1.55), fill=TINT, line=None)
    shape_text(c3, [
        [R("Our success bar", 13.5, NAVY, True)],
        [R("An order settles with no manual step — and every action "
           "is authorized and auditable.", 12.5, BODY, True)],
    ], space_after=5, m=0.2)
    notes(s, """
The brief asked for five core modules — Order Execution, Portfolio Management,
Reporting & Charting, Technical Analytics, Paper Trading — on a DevSecOps
foundation, using the program's simulation dataset. Our success bar: an order
settles with no manual step, and every action is authorized and auditable.
""")


def slide_05(prs):
    s = base_slide(prs, 5, "Platform overview — five modules + governance spine",
                   kicker="Section 2 · What we built")
    mods = [
        ("Order Execution & STP", "Ticket to settlement, zero manual steps"),
        ("Portfolio Management", "Positions, valuation, KPIs, live marks"),
        ("Reporting & Charting", "ECharts candles + indicators; PDF/CSV reports"),
        ("Technical Analytics", "SMA/EMA/RSI/MACD/Bollinger; price alerts"),
        ("Paper Trading", "Same pipeline as real trading — PAPER flag"),
    ]
    W = Inches(2.28)
    for i, (name, desc) in enumerate(mods):
        x = MARGIN + i * (W + Inches(0.19))
        c = card(s, x, Inches(1.5), W, Inches(1.62), fill=PANEL, line=HAIR)
        shape_text(c, [
            [R(name, 12, NAVY, True)],
            [R(desc, 10, MUTED)],
        ], space_after=5, m=0.14)
    spine = card(s, MARGIN, Inches(3.38), CONTENT_W, Inches(0.72),
                 fill=NAVY, line=None)
    shape_text(spine, [[
        R("GOVERNANCE SPINE   ", 11, WHITE, True),
        R("deny-by-default RBAC  ·  just-in-time access  ·  break-glass (4 h)  ·  "
          "hash-chained audit trail", 11.5, WHITE)]],
        align=PP_ALIGN.CENTER, m=0.1)
    stats = [
        "11 instruments — 7 dataset equities + 4 bonds the business asked for",
        "5 order types + time-in-force + trailing stops",
        "Price alerts with evaluator",
        "Scheduled PDF/CSV reports",
        "Paper trading on the same code path",
        "GenAI assistant — grounded, advisory-only",
    ]
    W2 = Inches(3.91)
    for i, t in enumerate(stats):
        x = MARGIN + (i % 3) * (W2 + Inches(0.2))
        y = Inches(4.42) + (i // 3) * Inches(1.02)
        c = card(s, x, y, W2, Inches(0.86), fill=WHITE, line=HAIR)
        shape_text(c, [[R(t, 11, BODY)]], m=0.14)
    notes(s, """
All five modules are live, plus the governance spine around them — RBAC,
just-in-time access, break-glass, and a hash-chained audit trail. Eleven
instruments: the seven dataset equities plus four bonds the business explicitly
asked for. Five order types plus time-in-force and trailing stops. Price
alerts, scheduled PDF/CSV reports, paper trading on the same pipeline as real
trading, and a GenAI assistant.
""")


def slide_06(prs):
    s = base_slide(prs, 6, "Architecture — one deployable, event-driven inside",
                   kicker="Section 2 · What we built")
    # React SPA
    spa = card(s, MARGIN, Inches(2.75), Inches(1.95), Inches(1.45), fill=WHITE,
               line=HAIR)
    shape_text(spa, [
        [R("React SPA", 13, NAVY, True)],
        [R("Vite · TypeScript", 9.5, MUTED)],
        [R("dark trading-terminal UI", 9.5, MUTED)],
    ], align=PP_ALIGN.CENTER, space_after=3, m=0.08)
    connect(s, Inches(2.55), Inches(3.47), Inches(3.28), Inches(3.47),
            color=ACCENT, w=1.75)
    add_text(s, Inches(2.1), Inches(3.08), Inches(1.6), Inches(0.3),
             [[R("REST · WS", 9, MUTED, True)]], align=PP_ALIGN.CENTER)
    # FastAPI container
    box = card(s, Inches(3.3), Inches(1.5), Inches(6.0), Inches(3.95),
               fill=PANEL, line=HAIR)
    shape_text(box, [[R("FastAPI modular monolith — 17 modules · one deployable",
                        12.5, NAVY, True)]],
               anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.CENTER, m=0.14)
    mods = ["orders", "portfolios", "marketdata", "+ 12 more"]
    for i, mname in enumerate(mods):
        chip = card(s, Inches(3.52) + i * Inches(1.44), Inches(2.1),
                    Inches(1.32), Inches(0.5), fill=WHITE, line=HAIR)
        shape_text(chip, [[R(mname, 9.5, BODY)]], align=PP_ALIGN.CENTER, m=0.03)
    pipeline = ["Transactional\noutbox", "Execution\nengine", "STP\nworker",
                "Settlement\nsweeper"]
    for i, name in enumerate(pipeline):
        x = Inches(3.52) + i * Inches(1.44)
        st = card(s, x, Inches(3.0), Inches(1.28), Inches(1.0), fill=TINT,
                  line=None)
        shape_text(st, [[R(part, 10, NAVY, True)]
                        for part in name.split("\n")],
                   align=PP_ALIGN.CENTER, space_after=0, m=0.04)
        if i < 3:
            connect(s, x + Inches(1.28), Inches(3.5), x + Inches(1.44),
                    Inches(3.5), color=MUTED, w=1.25)
    add_text(s, Inches(3.52), Inches(4.25), Inches(5.6), Inches(0.3),
             [[R("asynchronous · event-driven · idempotent consumers",
                 10, MUTED, False, True)]], align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.52), Inches(4.68), Inches(5.6), Inches(0.55),
             [[R("Outbox events written in the same DB transaction as the state "
                 "change; a relay publishes them to the bus.", 9.5, MUTED)]],
             align=PP_ALIGN.CENTER)
    # Data stores
    pg = card(s, Inches(9.85), Inches(1.85), Inches(2.88), Inches(1.25),
              fill=WHITE, line=HAIR)
    shape_text(pg, [
        [R("PostgreSQL / SQLite", 11.5, NAVY, True)],
        [R("transactions · audit · outbox", 9.5, MUTED)],
    ], align=PP_ALIGN.CENTER, space_after=3, m=0.08)
    rd = card(s, Inches(9.85), Inches(3.85), Inches(2.88), Inches(1.25),
              fill=WHITE, line=HAIR)
    shape_text(rd, [
        [R("Redis", 11.5, NAVY, True)],
        [R("sessions · event streams", 9.5, MUTED)],
    ], align=PP_ALIGN.CENTER, space_after=3, m=0.08)
    connect(s, Inches(9.3), Inches(2.47), Inches(9.85), Inches(2.47),
            color=MUTED, w=1.25)
    connect(s, Inches(9.3), Inches(4.47), Inches(9.85), Inches(4.47),
            color=MUTED, w=1.25)
    # Dataset strip + replay arrow
    ds = card(s, MARGIN, Inches(5.8), CONTENT_W, Inches(1.05), fill=TINT,
              line=None)
    shape_text(ds, [[
        R("SIMULATION DATASET   ", 10.5, NAVY, True),
        R("~190,000 price bars · 9,300 news items — replayed on a simulation "
          "clock (≈ one market day every 78 s).  ", 11.5, BODY),
        R("Nothing in the system knows the data is simulated.",
          11.5, NAVY, True, True)]], align=PP_ALIGN.CENTER, m=0.16)
    connect(s, Inches(6.3), Inches(5.8), Inches(6.3), Inches(5.45),
            color=ACCENT, w=1.5, dash=True)
    add_text(s, Inches(6.45), Inches(5.5), Inches(1.6), Inches(0.25),
             [[R("tick replay", 9, MUTED, True)]])
    notes(s, """
One diagram. A React single-page app talks REST and WebSocket to a FastAPI
modular monolith — sixteen modules, one deployable. Inside, an order flows
through a transactional outbox to the execution engine, the STP worker, and
the settlement sweeper — asynchronous, event-driven, and idempotent. Market
data is the program's dataset — about 190,000 price bars and 9,300 news items —
replayed on a *simulation clock*, so the platform runs in market time, roughly
a day every 78 seconds. Nothing in the system knows the data is simulated.
""")


def slide_07(prs):
    s = base_slide(prs, 7, "GenAI in the platform — grounded and advisory-only",
                   kicker="Section 2 · What we built")
    cards_ = [
        ("Grounded Q&A, with citations",
         "Answers questions grounded in your actual data — positions, KPIs, "
         "news — with citations. Refuses to invent figures it doesn't have."),
        ("News summarizer",
         "Summarizes real coverage for the selected instrument, with a "
         "sentiment score and headline citations."),
        ("Advisory-only guardrail",
         "It can suggest a trade — but the suggestion always lands in the "
         "standard order ticket: same validation, same two-click confirmation."),
    ]
    W = Inches(3.91)
    for i, (title, body) in enumerate(cards_):
        x = MARGIN + i * (W + Inches(0.2))
        c = card(s, x, Inches(1.55), W, Inches(3.15), fill=PANEL, line=HAIR)
        shape_text(c, [
            [R(title, 14, NAVY, True)],
            [R(body, 12, BODY)],
        ], anchor=MSO_ANCHOR.TOP, space_after=10, m=0.22)
    strip(s, Inches(5.15), Inches(1.05), [
        [R("Honest status: ", 12, NAVY, True),
         R("today a rule-based engine with a clean LLM seam — the interface is "
           "ready for a real model, and responses are honestly marked ",
          12, BODY),
         R("mock: true", 12, NAVY, True), R(".", 12, BODY)]])
    add_text(s, MARGIN, Inches(6.45), CONTENT_W, Inches(0.35),
             [[R("A wrong summary costs a read, not a trade.", 11.5, MUTED,
                 False, True)]], align=PP_ALIGN.CENTER)
    notes(s, """
Two honest things. First, the assistant answers questions *grounded in your
actual data* — positions, KPIs, news — with citations, and it refuses to invent
figures it doesn't have. Second, the news panel summarizes real coverage for
the selected instrument with a sentiment score and headline citations. And the
guardrail: the assistant is **advisory only** — it can suggest a trade, but the
suggestion always lands in the standard order ticket, with the same validation
and the same two-click confirmation as any other order. Today it's a rule-based
engine with a clean LLM seam — the interface is ready for a real model, and the
responses are honestly marked `mock: true`.
""")


# ============================================================= slides 8–14 ====
def slide_08(prs):
    s = base_slide(prs, 8, "Demo — 90 seconds, live on the simulation dataset",
                   kicker="Section 2 · What we built")
    add_text(s, MARGIN, Inches(1.35), CONTENT_W, Inches(0.3),
             [[R("Watch three things:", 12.5, BODY, True)]])
    beats = [
        ("1", "Buy 50 TSLA",
         "One click from the workspace — a confirmation card with the full "
         "cost impact — second click to submit."),
        ("2", "Live marks",
         "The position marks live over the WebSocket push — no refresh — "
         "and the risk panel reacts."),
        ("3", "SETTLED",
         "With no manual step, the settlement instruction transitions to "
         "settled — that is straight-through processing."),
    ]
    W = Inches(3.91)
    for i, (num, title, body) in enumerate(beats):
        x = MARGIN + i * (W + Inches(0.2))
        c = card(s, x, Inches(1.75), W, Inches(3.1), fill=PANEL, line=HAIR)
        shape_text(c, [
            [R(num, 26, ACCENT, True)],
            [R(title, 16, NAVY, True)],
            [R(body, 12, BODY)],
        ], anchor=MSO_ANCHOR.TOP, space_after=8, m=0.22)
        if i < 2:
            connect(s, x + W, Inches(3.3), x + W + Inches(0.2), Inches(3.3),
                    color=ACCENT, w=2.0)
    strip(s, Inches(5.3), Inches(1.0), [
        [R("Fallback: ", 12, NAVY, True),
         R("recorded video of the exact same path — we narrate the same three "
           "beats, say plainly the fallback is running, and offer to retry "
           "live during Q&A.", 12, BODY)]])
    add_text(s, MARGIN, Inches(6.5), CONTENT_W, Inches(0.3),
             [[R("Full checklist in script §10 — everything is localhost, no "
                 "venue Wi-Fi dependency.", 10.5, MUTED, False, True)]],
             align=PP_ALIGN.CENTER)
    notes(s, """
Ninety seconds, live. Watch three things: (1) I buy 50 TSLA from the workspace —
one click, a confirmation card with full cost impact, second click. (2) The
position marks live over the WebSocket push — no refresh — and the risk panel
reacts. (3) With no manual step, the settlement instruction transitions to
*settled* — that is straight-through processing. The settlement lifecycle is
visible in the UI: the settlements list shows every state, and ops can retry
exceptions.

Demo path (script §10): workspace → select TSLA → size 50 → BUY → confirmation
card → Confirm → fill toast → position marks live → risk panel reacts → order
status FILLED → settlement instruction → SETTLED in the settlements list.

Handover: "A fill is where the technology story ends — but it's where the
operations story begins. [P3]."
""")


def slide_09(prs):
    s = base_slide(prs, 9, "Operational processes — three workflows, one platform",
                   kicker="Section 3 · Operations")
    cols = [
        ("Trade settlement", [
            ("Execution creates a settlement instruction — ",
             "EXECUTED → AFFIRMED → SETTLED, automatically — every state "
             "visible in the settlements list.", 0),
            ("Exceptions are the job: ", "they surface on the operations "
             "dashboard — audited, and ops can retry them there.", 0),
            ("Same pipeline for client, house and paper books — ",
             "paper is the same code path with a portfolio-type flag.", 0),
        ]),
        ("Risk management", [
            ("Pre-trade validation on every order: ",
             "cash, lot size, configurable max notional, restricted-instrument "
             "list.", 0),
            ("Restricted orders reject with a reason — ",
             "and the rejection is audited.", 0),
            ("Post-trade, the risk panel computes ",
             "concentration, volatility and top holdings, live from "
             "positions.", 0),
        ]),
        ("System access", [
            ("Deny-by-default: ", "roles grant permissions; just-in-time "
             "grants expire automatically.", 0),
            ("Privileged actions via CyberArk-style checkout ",
             "(mocked adapter — the interface is real).", 0),
            ("Break-glass time-boxed to 4 h ", "with a 24 h review SLA.", 0),
            ("Hash-chained, append-only audit — ", "tamper-evident; the "
             "auditor role can search and export it.", 0),
        ]),
    ]
    W = Inches(3.91)
    for i, (title, items) in enumerate(cols):
        x = MARGIN + i * (W + Inches(0.2))
        c = card(s, x, Inches(1.5), W, Inches(4.75), fill=PANEL, line=HAIR)
        shape_text(c, [[R(title, 14, NAVY, True)]], anchor=MSO_ANCHOR.TOP,
                   m=0.18)
        add_bullets(s, x + Inches(0.18), Inches(2.1), W - Inches(0.36),
                    Inches(4.0), items, size=11, space_after=9)
    strip(s, Inches(6.45), Inches(0.55), [
        [R("The business workflow and the technical workflow are the same "
           "workflow.", 12.5, NAVY, True, True)]])
    notes(s, """
The corporate side of the team mapped three workflows, and the platform
implements all three.

Trade settlement: execution creates a settlement instruction that moves
EXECUTED → AFFIRMED → SETTLED automatically — every state visible in the
settlements list. Exceptions surface on the operations dashboard, with audit on
every transition, and ops can retry an exception from the same dashboard. The
same pipeline serves client, house and paper books — paper trading isn't a toy
mode, it's the same code path with a portfolio type flag.

Risk management: pre-trade, every order passes validation — cash, lot size, a
configurable max notional, and a restricted-instrument list managed by our
Security Administrator; restricted orders reject with a reason, and the
rejection is audited. Post-trade, the risk panel computes concentration,
volatility and top holdings live from positions.

System access: nobody gets anything by default. Roles grant permissions;
just-in-time grants expire automatically; privileged actions go through a
CyberArk-style checkout (mocked adapter — the interface is real, the vault is
the training stand-in); break-glass is time-boxed to four hours with a 24-hour
review SLA. Every denial, grant and activation is written to a hash-chained,
append-only audit trail — tamper-evident by construction, searchable and
exportable by the auditor role.

Handover: "Building that in three weeks took a way of working, not just code —
[P4]."
""")


def slide_10(prs):
    s = base_slide(prs, 10, "How we worked — three one-week sprints",
                   kicker="Section 4 · Process")
    weeks = [
        ("Week 1", "Walking skeleton + governance spine",
         "Auth · RBAC · audit · access-request lifecycle · dataset loader"),
        ("Week 2", "Trading core",
         "Order ticket · execution engine · STP settlement · portfolios · "
         "charts · notifications"),
        ("Week 3", "Breadth & hardening",
         "Paper trading · reports · analytics · admin dashboards · "
         "GenAI assistant"),
    ]
    W = Inches(3.85)
    for i, (wk, title, body) in enumerate(weeks):
        x = MARGIN + i * (W + Inches(0.29))
        shape = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        ch = card(s, x, Inches(1.5), W, Inches(0.62), fill=NAVY, line=None,
                  shape=shape)
        shape_text(ch, [[R(wk + " — " + title, 12.5, WHITE, True)]],
                   align=PP_ALIGN.CENTER, m=0.08)
        c = card(s, x, Inches(2.35), W, Inches(1.95), fill=PANEL, line=HAIR)
        shape_text(c, [[R(body, 11.5, BODY)]], anchor=MSO_ANCHOR.MIDDLE,
                   align=PP_ALIGN.CENTER, m=0.18)
        if i < 2:  # milestone feedback diamond between weeks
            d = card(s, x + W + Inches(0.035), Inches(1.7), Inches(0.22),
                     Inches(0.22), fill=ACCENT, line=None,
                     shape=MSO_SHAPE.DIAMOND)
    add_text(s, MARGIN, Inches(4.4), CONTENT_W, Inches(0.3),
             [[R("◆ milestone feedback after each sprint · daily facilitator "
                 "touchpoints", 10.5, MUTED)]], align=PP_ALIGN.CENTER)
    fb = card(s, MARGIN, Inches(4.95), CONTENT_W, Inches(1.8), fill=TINT,
              line=None)
    shape_text(fb, [
        [R("Feedback changed course. ", 13, NAVY, True),
         R("The product-owner feedback round after the trading demo produced "
           "four concrete changes:", 13, BODY)],
        [R("two-click confirmation  ·  bonds  ·  stop orders  ·  order "
           "restrictions", 13.5, ACCENT, True)],
        [R("Recorded in the repo as interview notes and a design document.",
           10.5, MUTED, False, True)],
    ], align=PP_ALIGN.CENTER, space_after=7, m=0.2)
    notes(s, """
Three one-week sprints. Week 1: walking skeleton plus the governance spine —
auth, RBAC, audit, the access-request lifecycle, and the dataset loader.
Week 2: the trading core — order ticket, execution engine, STP settlement,
portfolios, charts, notifications. Week 3: breadth and hardening — paper
trading, reports, analytics, admin dashboards, and the GenAI assistant. Daily
facilitator touchpoints; structured feedback after each milestone — and we
changed course because of it. The product-owner feedback round after the
trading demo produced four concrete changes you saw: the two-click
confirmation, bonds, stop orders, and order restrictions. That feedback is in
the repo as interview notes and a design document.
""")


def slide_11(prs):
    s = base_slide(prs, 11, "Engineering practices — verified, not asserted",
                   kicker="Section 4 · Process")
    add_bullets(s, MARGIN, Inches(1.5), Inches(6.55), Inches(4.7), [
        ("Design before code: ", "26 numbered design documents, one per "
         "module or major feature, each traced to SRS requirement IDs.", 0),
        ("Git flow: ", "feature branches off an integration branch — 13 "
         "merges, 12 feature branches in the history — plus a dated changelog "
         "entry per milestone.", 0),
        ("Verification gates on every milestone: ", "full test suite, strict "
         "TypeScript build, an end-to-end walk of the real stack, and "
         "headless-browser screenshots for UI changes.", 0),
        ("Second pair of eyes: ", "verified independently, never trusted on "
         "the author's word.", 0),
    ], size=12, space_after=12)
    add_text(s, Inches(7.5), Inches(1.42), Inches(5.2), Inches(0.3),
             [[R("SIX-STAGE GITLAB CI/CD", 10.5, MUTED, True)]])
    stages = [
        ("lint", "compileall + pip check"),
        ("test", "backend pytest · frontend build"),
        ("security scan", "gitleaks + trivy — blocking"),
        ("build", "Docker images for both tiers"),
        ("deploy-dev", "SSH · docker compose up -d --build"),
        ("deploy-demo", "manual gate"),
    ]
    y = Inches(1.78)
    for name, sub in stages:
        st = card(s, Inches(7.5), y, Inches(2.1), Inches(0.56), fill=PANEL,
                  line=HAIR)
        shape_text(st, [[R(name, 11, NAVY, True)]], align=PP_ALIGN.CENTER,
                   m=0.05)
        add_text(s, Inches(9.75), y + Inches(0.11), Inches(3.0), Inches(0.4),
                 [[R(sub, 9, MUTED)]])
        y += Inches(0.68)
    add_text(s, Inches(7.5), y + Inches(0.02), Inches(5.2), Inches(0.55),
             [[R("+ Docker images for both tiers · Terraform single-VM cloud "
                 "deployment (provider-portable).", 10.5, BODY)]])
    strip(s, Inches(6.35), Inches(0.8), [
        [R("Honest caveat: ", 11.5, NAVY, True),
         R("the pipeline, containers and Terraform are written and statically "
           "reviewed, but not executed — no cloud was available in the program "
           "environment. That stays on the roadmap.", 11.5, BODY)]])
    notes(s, """
Design before code: twenty-six numbered design documents, one per module or
major feature, each traced to SRS requirement IDs. Git flow: feature branches
off an integration branch — you can see thirteen merges and twelve feature
branches in the history — with a dated changelog entry per milestone. Every
milestone was verified the same way: the full test suite, the strict TypeScript
build, an end-to-end walk of the real stack, and headless-browser screenshots
for UI changes — checked by a second pair of eyes, never trusted on the
author's word. CI/CD is a defined six-stage GitLab pipeline — lint, test,
security scan, build, deploy-dev, deploy-demo — with Docker images and
Terraform for a single-VM cloud deployment. Honest caveat: written and
statically reviewed, but not executed — no cloud in the program environment.
That stays on the roadmap.
""")


def slide_12(prs):
    s = base_slide(prs, 12, "Cross-functional collaboration — with receipts",
                   kicker="Section 4 · Process")
    W = Inches(3.91)
    # Card 1 — division of labor
    c = card(s, MARGIN, Inches(1.5), W, Inches(4.55), fill=PANEL, line=HAIR)
    shape_text(c, [[R("Joint artifacts", 14, NAVY, True)]],
               anchor=MSO_ANCHOR.TOP, m=0.2)
    add_bullets(s, MARGIN + Inches(0.2), Inches(2.1), W - Inches(0.4),
                Inches(3.8), [
        ("Corporate owned ", "SRS traceability and the stakeholder "
         "interviews.", 0),
        ("18 open questions, ", "each assigned to a stakeholder.", 0),
        ("Technology owned ", "the build.", 0),
    ], size=11.5, space_after=9)
    # Card 2 — equities to bonds
    x2 = MARGIN + W + Inches(0.2)
    c = card(s, x2, Inches(1.5), W, Inches(4.55), fill=PANEL, line=HAIR)
    shape_text(c, [[R("Equities → bonds", 14, NAVY, True)]],
               anchor=MSO_ANCHOR.TOP, m=0.2)
    add_bullets(s, x2 + Inches(0.2), Inches(2.1), W - Inches(0.4),
                Inches(3.8), [
        ("The product-owner interview: ", "\u201cequities only is not "
         "credible.\u201d", 0),
        ("Two days later ", "the platform priced bonds properly.", 0),
        ("Quoted as percent of par, ", "with correct cash math.", 0),
    ], size=11.5, space_after=9)
    # Card 3 — the single-click misalignment
    x3 = x2 + W + Inches(0.2)
    c = card(s, x3, Inches(1.5), W, Inches(4.55), fill=PANEL, line=HAIR)
    shape_text(c, [[R("The \u201csingle-click\u201d misalignment", 14, NAVY, True)]],
               anchor=MSO_ANCHOR.TOP, m=0.2)
    add_bullets(s, x3 + Inches(0.2), Inches(2.1), W - Inches(0.4),
                Inches(1.5), [
        ("Technology heard ", "one click.", 0),
        ("The business meant ", "one informed decision.", 0),
    ], size=11.5, space_after=9)
    flow = ["click to arm", "confirmation card — full cost impact",
            "click to submit"]
    fw = [Inches(1.0), Inches(1.55), Inches(1.0)]
    fx = x3 + Inches(0.18)
    for j, label in enumerate(flow):
        ch = card(s, fx, Inches(3.6), fw[j], Inches(0.85), fill=TINT, line=None,
                  shape=MSO_SHAPE.CHEVRON)
        shape_text(ch, [[R(label, 9, NAVY, True)]], align=PP_ALIGN.CENTER,
                   m=0.05)
        fx += fw[j] + Inches(0.06)
    add_text(s, x3 + Inches(0.2), Inches(4.75), W - Inches(0.4), Inches(1.1),
             [[R("The resolution was the two-click flow:", 11, BODY, True)],
              [R("\u201cFaster than a ticket, safer than a blind click.\u201d",
                 11.5, ACCENT, False, True)]], space_after=4)
    strip(s, Inches(6.35), Inches(0.6), [
        [R("Day to day: corporate analysts owned requirements and interviews; "
           "technology owned the build — one team, one backlog.", 11.5, BODY)]])
    notes(s, """
Day to day this meant: corporate analysts owned the SRS traceability and the
stakeholder interviews — eighteen open questions, each assigned to a
stakeholder — and technology owned the build. The product-owner interview is
the clearest joint artifact: the business said "equities only is not credible,"
and two days later the platform priced bonds properly, quoted as percent of
par with correct cash math. One real misalignment: "single-click trading" —
technology heard *one* click, the business meant *one informed decision*. The
resolution was the two-click flow: one click to arm, a confirmation card with
the full cost impact, one click to submit. Faster than a ticket, safer than a
blind click.

Handover: "That's how we worked. What actually held up — and what didn't —
[P1]."
""")


def slide_13(prs):
    s = base_slide(prs, 13, "Challenges & solutions — three war stories",
                   kicker="Section 5 · Retrospective")
    stories = [
        ("The dataset that silently didn't load",
         "After switching to the real dataset, every stock showed \u201cNo "
         "price data.\u201d Legacy rows from an old dev database made the "
         "loader's global \u201cis the table empty?\u201d check skip the "
         "entire tick load.",
         "Found by inspecting the database, not by guessing; load per symbol, "
         "with a regression test."),
        ("The test that only failed after midnight",
         "A news fixture used relative-to-now timestamps; between midnight and "
         "2 a.m. UTC the two items straddled a day boundary and the suite went "
         "red. It only ever passed because of the hour it ran.",
         "Anchored deterministically — and we re-audited every time-sensitive "
         "test."),
        ("Real-time wasn't a library import",
         "Cancellations mid-database-call could wedge the event loop, and "
         "naive datetimes compared *silently* wrong against timezone-aware "
         "ones.",
         "Both are now documented pitfalls, with fixes and tests behind them."),
    ]
    y = Inches(1.45)
    for title, problem, fix in stories:
        band = card(s, MARGIN, y, CONTENT_W, Inches(1.62), fill=PANEL,
                    line=HAIR)
        shape_text(band, [[R(title, 12.5, NAVY, True)]],
                   anchor=MSO_ANCHOR.TOP, m=0.16)
        add_text(s, MARGIN + Inches(0.18), y + Inches(0.52), Inches(7.5),
                 Inches(1.05), [[R(problem, 10.5, BODY)]], line_spacing=1.05)
        fc = card(s, MARGIN + Inches(7.95), y + Inches(0.14), Inches(4.0),
                  Inches(1.34), fill=TINT, line=None)
        shape_text(fc, [[R("Fix — ", 10.5, NAVY, True), R(fix, 10.5, BODY)]],
                   m=0.14, line_spacing=1.05)
        y += Inches(1.82)
    notes(s, """
One story live — the dataset one, in forty seconds. The other two stay on the
slide and in these notes as Q&A ammunition (script §9) for any "what was hard /
what went wrong?" question.

The dataset that silently didn't load: after switching to the real dataset,
every stock showed "No price data." Root cause: legacy rows from an old dev
database made the loader's global "is the table empty?" check skip the entire
tick load. Found by inspecting the database, not by guessing; fixed by loading
per symbol, with a regression test.

The test that only failed after midnight: a news fixture used relative-to-now
timestamps; between midnight and 2 a.m. UTC the two items straddled a day
boundary and the suite went red. It only ever passed because of the hour it
ran. Now anchored deterministically — and it made us re-audit every
time-sensitive test.

Real-time wasn't a library import: cancellations mid-database-call could wedge
the event loop, and naive datetimes compared *silently* wrong against
timezone-aware ones. Both are now documented pitfalls with fixes and tests
behind them.
""")


def slide_14(prs):
    s = base_slide(prs, 14, "What worked / what didn't", kicker="Section 5 · Retrospective")
    W = Inches(5.95)
    c1 = card(s, MARGIN, Inches(1.5), W, Inches(5.15), fill=PANEL, line=HAIR)
    shape_text(c1, [[R("What worked", 15, ACCENT, True)]],
               anchor=MSO_ANCHOR.TOP, m=0.2)
    add_bullets(s, MARGIN + Inches(0.2), Inches(2.15), W - Inches(0.4),
                Inches(4.3), [
        ("Design-first: ", "26 design docs meant feedback rounds changed "
         "documents before they changed code.", 0),
        ("The event pipeline: ", "the STP flow never needed a manual hack.", 0),
        ("Verification discipline: ", "~100 tests, a strict build, and "
         "screenshot-verified UI rounds.", 0),
    ], size=12, space_after=12)
    x2 = MARGIN + W + Inches(0.23)
    c2 = card(s, x2, Inches(1.5), W, Inches(5.15), fill=PANEL, line=HAIR)
    shape_text(c2, [[R("What didn't", 15, MUTED, True)]],
               anchor=MSO_ANCHOR.TOP, m=0.2)
    add_bullets(s, x2 + Inches(0.2), Inches(2.15), W - Inches(0.4),
                Inches(4.3), [
        ("No frontend tests: ", "UI quality rests on build strictness and "
         "manual verification.", 0),
        ("Deployment stack unexecuted: ", "no cloud in the program "
         "environment.", 0),
        ("GenAI reality check: ", "a rule-based assistant is robust and "
         "honest, but it is not a language model — the summary-quality "
         "ceiling is real, and that's a phase-2 decision, not a bug.", 0),
    ], size=12, space_after=12)
    notes(s, """
Worked: design-first — 26 design docs meant feedback rounds changed documents
before they changed code. The event pipeline — the STP flow never needed a
manual hack. Verification discipline — ~100 tests, a strict build, and
screenshot-verified UI rounds.

What didn't: we have **no frontend tests** — UI quality rests on build
strictness and manual verification. The deployment stack is unexecuted — no
cloud in the program environment. And the GenAI reality check: a rule-based
assistant is robust and honest, but it is not a language model — the summary
quality ceiling is real, and that's a phase-2 decision, not a bug.

Handover: "So what did three weeks teach us, and where does it go next —
[P3]."
""")


# ============================================================ slides 15–18 ====
def slide_15(prs):
    s = base_slide(prs, 15, "Key learnings", kicker="Section 6 · Learnings")
    cols = [
        ("Technical", [
            ("Event-driven design pays for itself ", "the first time a "
             "requirement changes mid-project — trailing stops and report "
             "scheduling slotted into the existing pipeline.", 0),
            ("GenAI in production ", "is a guardrail problem before it's a "
             "model problem.", 0),
        ]),
        ("Domain", [
            ("Settlement is a workflow, not a status field — ", "the "
             "corporate mapping of affirm/settle semantics made the STP demo "
             "credible.", 0),
            ("A bond is not a stock with a different symbol — ",
             "percent-of-par quoting, yield and duration had to be "
             "first-class.", 0),
        ]),
        ("Professional", [
            ("The stakeholder-voices framing changed how we built — ",
             "every feature traces to a person, and prioritization arguments "
             "got short.", 0),
            ("Honest status beats impressive status — ", "facilitators know "
             "what three weeks allows.", 0),
        ]),
    ]
    W = Inches(3.91)
    for i, (title, items) in enumerate(cols):
        x = MARGIN + i * (W + Inches(0.2))
        c = card(s, x, Inches(1.5), W, Inches(4.9), fill=PANEL, line=HAIR)
        shape_text(c, [[R(title, 14, NAVY, True)]], anchor=MSO_ANCHOR.TOP,
                   m=0.2)
        add_bullets(s, x + Inches(0.2), Inches(2.15), W - Inches(0.4),
                    Inches(4.0), items, size=11.5, space_after=11)
    notes(s, """
(Compressed for the 10-minute slot: one line per category — the full text
below is reference.)

Technical: event-driven design pays for itself the first time a requirement
changes mid-project — trailing stops and report scheduling both slotted into
the existing pipeline. And GenAI in production is a guardrail problem before
it's a model problem.

Domain: settlement is a workflow, not a status field — the corporate mapping
of affirm/settle semantics is what made the STP demo credible. And "a bond is
not a stock with a different symbol" — percent-of-par quoting, yield and
duration all had to be first-class.

Professional: the stakeholder-voices framing changed how we built — every
feature traces to a person, and that made prioritization arguments short. And
honest status beats impressive status: facilitators know what three weeks
allows.

Transition (same speaker continues): "Which brings us back to the four people
we started with."
""")


def slide_16(prs):
    s = base_slide(prs, 16, "Metrics — honest numbers", kicker="Section 6 · Roadmap & close")
    tiles = [
        ("42", "commits on the integration branch", False),
        ("14", "merges", False),
        ("12", "feature branches", False),
        ("~100", "backend tests — all green", False),
        ("17", "backend modules", False),
        ("26", "design documents — SRS-traced", False),
        ("11", "instruments — 7 equities + 4 bonds", False),
        ("~190k", "price bars replayed on the sim clock", False),
        ("9.3k", "news items as reference data", False),
        ("0", "pipeline runs — defined & reviewed, not executed", True),
    ]
    W, H = Inches(2.33), Inches(2.0)
    for i, (num, label, flag) in enumerate(tiles):
        x = MARGIN + (i % 5) * (W + Inches(0.12))
        y = Inches(1.6) + (i // 5) * (H + Inches(0.2))
        c = card(s, x, y, W, H, fill=(TINT if flag else PANEL),
                 line=(ACCENT if flag else HAIR))
        shape_text(c, [
            [R(num, 30, (ACCENT if flag else NAVY), True)],
            [R(label, 9.5, BODY)],
        ], align=PP_ALIGN.CENTER, space_after=4, m=0.1)
    strip(s, Inches(6.1), Inches(0.7), [
        [R("That zero is on the roadmap — ", 12, NAVY, True),
         R("the CI/CD is defined and reviewed, not yet executed (no cloud in "
           "the program environment).", 12, BODY)]])
    notes(s, """
42 commits on the integration branch, 14 merges, 12 feature branches. ~100
backend tests, all green. 17 backend modules, 26 design documents, a dated
changelog for every milestone. 11 instruments, ~190k price bars and 9.3k news
items replayed on the simulation clock. Zero pipeline runs — the CI/CD is
defined and reviewed, not yet executed. That zero is on the roadmap.
""")


def slide_17(prs):
    s = base_slide(prs, 17, "Future roadmap — answering the four voices",
                   kicker="Section 6 · Roadmap & close")
    quads = [
        ("Roy", "scalability & DevSecOps", [
            "Run the pipeline for real; Alembic migrations",
            "Multi-instance via the Redis session / event-bus path that already exists",
            "External audit anchoring",
        ]),
        ("Tom & Patricia", "adoption", [
            "Keep the dashboard the single screen they asked for",
            "Live news via the provider seam the business already reviewed",
            "Change management is a feature, not an afterthought",
        ]),
        ("Nora", "maintainability", [
            "The assistant's LLM seam lets a real model in without touching the guardrails",
            "Batch-system knowledge preserved in 26 design docs and a runbook",
            "Built with her skepticism, not against it",
        ]),
        ("Rohan", "cost/benefit for the CFO", [
            "Iceberg orders and per-desk limits: designed, deferred",
            "Roadmap sequenced by business value",
            "The demo you saw runs on a laptop",
        ]),
    ]
    W, H = Inches(5.95), Inches(2.5)
    for i, (name, theme, points) in enumerate(quads):
        x = MARGIN + (i % 2) * (W + Inches(0.23))
        y = Inches(1.5) + (i // 2) * (H + Inches(0.22))
        c = card(s, x, y, W, H, fill=PANEL, line=HAIR)
        shape_text(c, [[
            R(name + "  ", 14, NAVY, True), R("— " + theme, 10.5, MUTED)]],
            anchor=MSO_ANCHOR.TOP, m=0.2)
        add_bullets(s, x + Inches(0.2), y + Inches(0.62), W - Inches(0.4),
                    H - Inches(0.8), [(p, "", 0) for p in points],
                    size=11, space_after=6)
    notes(s, """
Roy — scalability & DevSecOps: run the pipeline for real; Alembic migrations;
multi-instance deployment with the Redis session and event-bus path that
already exists; external audit anchoring.

Tom & Patricia — adoption: keep the dashboard the single screen they asked
for; live news via the provider seam the business already reviewed; change
management is a feature, not an afterthought.

Nora — maintainability: the assistant's LLM seam lets a real model in without
touching the guardrails; the batch system's knowledge is preserved in 26
design docs and a runbook — we built *with* her skepticism, not against it.

Rohan — cost/benefit for the CFO: iceberg orders and per-desk limits are
designed but deferred; the roadmap is sequenced by business value, and the
demo you saw runs on a laptop.
""")


def slide_18(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.09))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    _no_shadow(bar)
    add_text(s, MARGIN, Inches(0.9), CONTENT_W, Inches(0.9),
             [[R("Thank you — questions?", 36, NAVY, True)]],
             align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(1.95), CONTENT_W, Inches(0.35),
             [[R("We set out to answer four people:", 13, MUTED, False, True)]],
             align=PP_ALIGN.CENTER)
    recap = [
        ("Rohan", "one-click trading that's still an informed decision"),
        ("Tom & Patricia", "a dashboard that replaces Excel without teaching a new system"),
        ("Nora", "a rebuild that respects what the batch system knows"),
        ("Roy", "DevSecOps that was day-one, not day-twenty"),
    ]
    W = Inches(8.6)
    x = (SLIDE_W - W) / 2
    y = Inches(2.6)
    for name, line in recap:
        c = card(s, x, y, W, Inches(0.62), fill=PANEL, line=HAIR)
        shape_text(c, [[R(name + "   ", 12.5, NAVY, True),
                        R("— " + line, 12.5, BODY)]], m=0.2)
        y += Inches(0.78)
    add_text(s, MARGIN, Inches(6.15), CONTENT_W, Inches(0.35),
             [[R("[Team name]  ·  Nomura Tech Graduate Program 2026  ·  "
                 "30 July 2026", 11.5, MUTED)]], align=PP_ALIGN.CENTER)
    add_text(s, MARGIN, Inches(6.55), CONTENT_W, Inches(0.3),
             [[R("Appendix: A1 trade-flow sequence · A2 data model · "
                 "A3 CI/CD pipeline", 10.5, MUTED, False, True)]],
             align=PP_ALIGN.CENTER)
    notes(s, """
We set out to answer four people: one-click trading that's still an informed
decision; a dashboard that replaces Excel without teaching a new system; a
rebuild that respects what the batch system knows; and DevSecOps that was
day-one, not day-twenty. Thank you — questions?

Q&A ammo in script §9 (10 likely questions: ROI of rebuilding, why not keep
batch, passwordless login, GenAI wrongness, no live market data, scaling,
partial fills, what runs in CI, looping sim clock, one more week). Appendix
slides A1–A3 on standby.
""")


# ============================================================== appendix =====
def slide_a1(prs):
    s = base_slide(prs, "A1", "Appendix — STP trade-flow sequence",
                   kicker="Q&A backup")
    stages = [
        ("Trader — workspace", "two-click order ticket"),
        ("Orders API", "validate: cash · lot size · max notional · restricted list"),
        ("Execution engine", "matches order vs live tick stream"),
        ("STP worker", "execution → position · cash · settlement instruction"),
        ("Settlement sweeper", "EXECUTED → AFFIRMED → SETTLED"),
    ]
    W, H = Inches(2.28), Inches(1.35)
    y = Inches(1.55)
    for i, (title, sub) in enumerate(stages):
        x = MARGIN + i * (W + Inches(0.19))
        c = card(s, x, y, W, H, fill=(TINT if i == 4 else PANEL),
                 line=(None if i == 4 else HAIR))
        shape_text(c, [
            [R(title, 11.5, NAVY, True)],
            [R(sub, 9.5, BODY)],
        ], align=PP_ALIGN.CENTER, space_after=4, m=0.1)
        if i < 4:
            connect(s, x + W, y + H / 2, x + W + Inches(0.19), y + H / 2,
                    color=ACCENT, w=1.75)
    # settlement-state chevrons
    add_text(s, MARGIN, Inches(3.35), CONTENT_W, Inches(0.28),
             [[R("SETTLEMENT STATES — AUTOMATIC, AUDITED AT EVERY TRANSITION",
                 10, MUTED, True)]])
    states = ["EXECUTED", "AFFIRMED", "SETTLED"]
    for i, st in enumerate(states):
        x = MARGIN + i * Inches(2.15)
        ch = card(s, x, Inches(3.72), Inches(1.95), Inches(0.55),
                  fill=NAVY, line=None,
                  shape=(MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON))
        shape_text(ch, [[R(st, 11.5, WHITE, True)]], align=PP_ALIGN.CENTER,
                   m=0.05)
    add_text(s, MARGIN + Inches(6.7), Inches(3.78), Inches(5.4), Inches(0.45),
             [[R("No manual step anywhere on the happy path.", 11, NAVY, True,
                 True)]])
    # supporting notes
    add_bullets(s, MARGIN, Inches(4.75), CONTENT_W, Inches(2.1), [
        ("Validation failures reject with a reason — ", "and the rejection "
         "is audited.", 0),
        ("State change + outbox event commit in the same transaction — ",
         "a relay publishes to the bus; consumers are idempotent.", 0),
        ("Exceptions are the job: ", "they surface on the operations "
         "dashboard for the Operations Analyst.", 0),
        ("Whole fills only in MVP — ", "executions are separate records, so "
         "partial fills need no schema change later.", 0),
    ], size=11.5, space_after=8)
    notes(s, """
If asked about the trade flow: order → validation (cash, lot size, max
notional, restricted list) → execution against the tick stream → STP worker
(position/cash + settlement instruction) → settlement sweeper moves EXECUTED →
AFFIRMED → SETTLED automatically, audit on every transition.

If asked about partial fills (script §9 Q7): out of MVP scope by design —
orders fill whole or rest working; executions are separate records, so
partials are admitted without a schema change.

If asked about exceptions: they surface on the operations dashboard with the
full audit trail.
""")


def slide_a2(prs):
    s = base_slide(prs, "A2", "Appendix — data model overview",
                   kicker="Q&A backup")
    groups = [
        ("Identity & access", ["users", "roles", "permissions", "grants",
                               "access requests", "break-glass"]),
        ("Market data", ["instruments", "price ticks", "news items",
                         "news sentiment"]),
        ("Trading", ["portfolios", "orders", "executions",
                     "settlement instructions", "positions"]),
        ("Assurance", ["valuation snapshots", "audit events", "outbox events",
                       "reports"]),
    ]
    W, H = Inches(5.95), Inches(1.98)
    for i, (title, entities) in enumerate(groups):
        x = MARGIN + (i % 2) * (W + Inches(0.23))
        y = Inches(1.5) + (i // 2) * (H + Inches(0.2))
        c = card(s, x, y, W, H, fill=PANEL, line=HAIR)
        shape_text(c, [[R(title, 12.5, NAVY, True)]], anchor=MSO_ANCHOR.TOP,
                   m=0.16)
        cx, cy = x + Inches(0.16), y + Inches(0.52)
        ch_h = Inches(0.36)
        for e in entities:
            cw = Inches(0.34 + 0.082 * len(e))
            if cx + cw > x + W - Inches(0.16):
                cx = x + Inches(0.16)
                cy += ch_h + Inches(0.12)
            chip = card(s, cx, cy, cw, ch_h, fill=WHITE, line=HAIR, radius=0.5)
            shape_text(chip, [[R(e, 9.5, BODY)]], align=PP_ALIGN.CENTER,
                       m=0.02)
            cx += cw + Inches(0.12)
    chain = card(s, MARGIN, Inches(5.9), CONTENT_W, Inches(1.0), fill=TINT,
                 line=None)
    shape_text(chain, [
        [R("AUDIT HASH CHAIN   ", 10.5, NAVY, True),
         R("payload_hash = sha256(canonical_json(entry) + prev_hash)", 12,
           BODY, True)],
        [R("Append-only, tamper-evident by construction — every denial, grant "
           "and activation is a link in the chain; the auditor role can "
           "search and export it.", 10.5, BODY)],
    ], align=PP_ALIGN.CENTER, space_after=4, m=0.16)
    notes(s, """
If asked about the data model: four clusters — identity & access (users,
roles, permissions, grants, access requests, break-glass), market data
(instruments, price ticks, news + sentiment), trading (portfolios, orders,
executions, settlement instructions, positions), assurance (valuation
snapshots, audit events, outbox events, reports).

If asked "your login is passwordless — how is that security?" (script §9 Q3):
it isn't, and we say so — dev-login is a training-only flag. The real model is
server-side sessions, deny-by-default RBAC re-checked per request,
hash-chained audit; SSO via OIDC is the designed integration point.

If asked how it scales (Q6): stateless API; Redis sessions and event streams
already implemented behind config flags; module boundaries drawn for a
mechanical split.
""")


def slide_a3(prs):
    s = base_slide(prs, "A3", "Appendix — CI/CD pipeline configuration",
                   kicker="Q&A backup")
    stages = [
        ("1 · lint", "compileall +\npip check", False),
        ("2 · test", "backend pytest\nfrontend build", False),
        ("3 · scan", "gitleaks + trivy\nHIGH/CRITICAL", True),
        ("4 · build", "Docker images\n(default branch)", False),
        ("5 · deploy-dev", "SSH: git pull +\ncompose up --build", False),
        ("6 · deploy-demo", "manual gate", False),
    ]
    W, H = Inches(1.88), Inches(1.3)
    y = Inches(1.7)
    for i, (name, sub, blocking) in enumerate(stages):
        x = MARGIN + i * (W + Inches(0.17))
        c = card(s, x, y, W, H, fill=(TINT if blocking else PANEL),
                 line=(ACCENT if blocking else HAIR))
        shape_text(c, [[R(name, 11.5, NAVY, True)]] +
                   [[R(part, 9, BODY)] for part in sub.split("\n")],
                   align=PP_ALIGN.CENTER, space_after=2, m=0.06)
        if i < 5:
            connect(s, x + W, y + H / 2, x + W + Inches(0.17), y + H / 2,
                    color=MUTED, w=1.5)
    add_text(s, MARGIN, Inches(3.2), CONTENT_W, Inches(0.3),
             [[R("Security scan BLOCKS the pipeline — secrets (gitleaks) and "
                 "HIGH/CRITICAL findings (trivy) fail the build.",
                 11, ACCENT, True)]], align=PP_ALIGN.CENTER)
    add_bullets(s, MARGIN, Inches(3.85), CONTENT_W, Inches(1.9), [
        ("Defined in .gitlab-ci.yml; ", "runs on merge requests and the "
         "default branch.", 0),
        ("Docker: ", "backend python:3.13-slim + uvicorn; frontend "
         "node:22-alpine build → nginx serving the SPA.", 0),
        ("Terraform: ", "single-VM reference deployment, optional RDS "
         "PostgreSQL, provider-portable (Azure parity notes).", 0),
    ], size=11.5, space_after=8)
    strip(s, Inches(6.0), Inches(0.85), [
        [R("Honest status: ", 11.5, NAVY, True),
         R("written and statically reviewed — zero pipeline runs (no cloud in "
           "the program environment). Gates were replicated locally on every "
           "merge.", 11.5, BODY)]])
    notes(s, """
If asked "you claim DevOps — what actually runs in CI?" (script §9 Q8): today,
lint, tests, and security scans are defined to block merges; the pipeline file
is real. What we did not do in the program environment is execute it in
GitLab — no cloud access — so we replicated the gates locally on every merge.

If asked "what would you do with one more week?" (Q10): run the deployment
stack for real, add the load test we specified (200 concurrent push clients),
and frontend tests — in that order.
""")


# ================================================================= build =====
def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for fn in (slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
               slide_07, slide_08, slide_09, slide_10, slide_11, slide_12,
               slide_13, slide_14, slide_15, slide_16, slide_17, slide_18,
               slide_a1, slide_a2, slide_a3):
        fn(prs)
    prs.save(OUT)
    return OUT


def verify(path: Path) -> bool:
    """Re-open the deck and check slide count, notes coverage, empty frames."""
    prs = Presentation(path)
    slides = list(prs.slides)
    ok = True
    print(f"\nVerification of {path.name}")
    print(f"  file size: {path.stat().st_size / 1024:.0f} KiB")
    print(f"  slide count: {len(slides)} (expected 21)")
    if len(slides) != 21:
        ok = False
        print("  FAIL: unexpected slide count")
    min_dim = Inches(0.25)  # smaller shapes are decorative (bars, diamonds)
    tol = Inches(0.02)      # rounding slack for the off-canvas bounds check
    for i, slide in enumerate(slides, 1):
        n_shapes = len(slide.shapes)
        has_notes = (slide.has_notes_slide
                     and bool(slide.notes_slide.notes_text_frame.text.strip()))
        if not has_notes:
            ok = False
        empty = []
        offcanvas = []
        for sh in slide.shapes:
            # bounding box (connector widths/heights can be negative)
            x0, x1 = sorted((sh.left, sh.left + sh.width))
            y0, y1 = sorted((sh.top, sh.top + sh.height))
            if x0 < -tol or y0 < -tol or x1 > SLIDE_W + tol or y1 > SLIDE_H + tol:
                offcanvas.append(sh.shape_id)
                ok = False
            if not getattr(sh, "has_text_frame", False):
                continue
            if sh.width < min_dim or sh.height < min_dim:
                continue  # decorative chrome (accent bars, rules, diamonds)
            if not sh.text_frame.text.strip():
                empty.append(sh.shape_id)
        if empty:
            ok = False
        flag = ""
        if not has_notes:
            flag += "  [NO NOTES]"
        if empty:
            flag += f"  [EMPTY TEXT shapes: {empty}]"
        if offcanvas:
            flag += f"  [OFF-CANVAS shapes: {offcanvas}]"
        print(f"  slide {i:>2}: shapes={n_shapes:>2} notes={'yes' if has_notes else 'NO '}"
              f"{flag}")
    print("  RESULT:", "PASS" if ok else "FAIL")
    return ok


if __name__ == "__main__":
    out = build()
    print(f"Wrote {out}")
    raise SystemExit(0 if verify(out) else 1)
